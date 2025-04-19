from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Helper function to convert Western Arabic numerals to Persian numerals
def to_persian_number(number):
    persian_digits = {
        '0': '۰',
        '1': '۱',
        '2': '۲',
        '3': '۳',
        '4': '۴',
        '5': '۵',
        '6': '۶',
        '7': '۷',
        '8': '۸',
        '9': '۹'
    }
    return ''.join(persian_digits[digit] for digit in str(number))

def add_wrapping_progress_circles(pptx_path, output_path):
    prs = Presentation(pptx_path)
    num_slides = len(prs.slides)

    # Slide dimensions
    width = prs.slide_width
    height = prs.slide_height

    # Circle settings
    diameter = Pt(30)
    spacing = Pt(6)
    margin = Pt(10)
    step = diameter + spacing

    # Adjust max positions to avoid corner overlap
    max_vertical = int((height - 2 * margin - diameter) // step) + 1  # Account for corner space
    max_horizontal = int((width - 2 * margin - diameter) // step) + 1  # Account for corner space
    max_positions = 2 * max_vertical + 2 * max_horizontal - 4  # Subtract 4 to avoid corner overlap

    if num_slides > max_positions:
        raise ValueError("Too many slides to fit progress indicators on slide edges.")

    # Generate edge positions: right → bottom → left → top
    positions = []

    # Right edge (top to bottom, excluding bottom corner)
    for i in range(max_vertical):
        x = width - margin - diameter
        y = margin + i * step
        positions.append((x, y))

    # Bottom edge (right to left, excluding both corners)
    for i in range(1, max_horizontal - 1):
        x = width - margin - diameter - i * step
        y = height - margin - diameter
        positions.append((x, y))

    # Left edge (bottom to top, excluding top corner)
    for i in range(max_vertical - 1, -1, -1):
        x = margin
        y = margin + i * step
        positions.append((x, y))

    # Top edge (left to right, excluding both corners)
    for i in range(1, max_horizontal - 1):
        x = margin + i * step
        y = margin
        positions.append((x, y))

    # Add indicators to each slide
    for slide_index, slide in enumerate(prs.slides):
        for i in range(num_slides):
            x, y = positions[i]
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, int(x), int(y), int(diameter), int(diameter)
            )
            fill = shape.fill
            fill.solid()
            shape.line.fill.background()

            if i < slide_index:
                fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow for seen
            elif i == slide_index:
                fill.fore_color.rgb = RGBColor(255, 0, 0)    # Red for current
            else:
                fill.fore_color.rgb = RGBColor(200, 200, 200)  # Gray for not seen

            # Add number to:
            # - the current slide's circle (red one)
            # - the last circle (with total number of slides)
            if i == slide_index or i == num_slides - 1:
                text_frame = shape.text_frame
                if i == slide_index:
                    text_frame.text = to_persian_number(i + 1)  # Persian numerals
                elif i == num_slides - 1:
                    text_frame.text = to_persian_number(num_slides)  # Persian numerals

                p = text_frame.paragraphs[0]
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = 1  # Center
                text_frame.word_wrap = False
                text_frame.auto_size = None
                text_frame.margin_top = 0
                text_frame.margin_bottom = 0
                text_frame.margin_left = 0
                text_frame.margin_right = 0

    last_slide = prs.slides[-1]
    textbox = last_slide.shapes.add_textbox(
        left=Pt(30), top=height - Pt(80), width=width - Pt(60), height=Pt(40)
    )
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Total Slides: {num_slides}"  # Unchanged, English style
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    prs.save(output_path)
    print(f"Progress indicators and total count added to: {output_path}")

# Example usage
# add_wrapping_progress_circles("input.pptx", "output.pptx")
name = 'Title Lorem Ipsum'
add_wrapping_progress_circles(f'{name}.pptx', f'{name}_edit.pptx')