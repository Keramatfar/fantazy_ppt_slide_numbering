from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

# Helper function to convert Western Arabic numerals to Persian numerals
def to_persian_number(number):
    persian_digits = {
        '0': '۰',
        '1': '۱',
        '2': '۲',
        '3': '۳',
        '4': '۴',
        '5': '۵',
        '6': '۶',
        '7': '۷',
        '8': '۸',
        '9': '۹'
    }
    return ''.join(persian_digits[digit] for digit in str(number))

def add_slide_numbers(pptx_path, output_path):
    # Load the presentation
    prs = Presentation(pptx_path)
    num_slides = len(prs.slides)

    # Add slide number to each slide
    for slide_index, slide in enumerate(prs.slides):
        # Slide dimensions
        width = prs.slide_width
        height = prs.slide_height

        # Add textbox at bottom-left
        left = Pt(30)
        top = height - Pt(50)  # Position near bottom
        textbox_width = Pt(100)
        textbox_height = Pt(20)

        textbox = slide.shapes.add_textbox(
            left=left,
            top=top,
            width=textbox_width,
            height=textbox_height
        )

        # Add slide number in format "current/total" with Persian numerals
        current_slide = to_persian_number(slide_index + 1)
        total_slides = to_persian_number(num_slides)
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{current_slide}/{total_slides}"  # e.g., "۲/۵۸"

        # Format the text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = "B Mitra"  # Set font to B Mitra
        p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    # Save the updated presentation
    prs.save(output_path)
    print(f"Slide numbers with B Mitra font added to: {output_path}")

add_slide_numbers(f"{name}.pptx", f"{name}_simple.pptx")