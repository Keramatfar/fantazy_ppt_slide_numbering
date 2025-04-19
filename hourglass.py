from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from pptx.dml.color import RGBColor

def add_wrapping_progress_circles(pptx_path, output_path):
    prs = Presentation(pptx_path)
    num_slides = len(prs.slides)

    # Slide dimensions
    width = prs.slide_width
    height = prs.slide_height

    # Circle settings
    diameter = Pt(30)  # slightly larger for text
    spacing = Pt(6)
    margin = Pt(10)
    step = diameter + spacing

    max_vertical = int((height - 2 * margin) // step)
    max_horizontal = int((width - 2 * margin) // step)
    max_positions = 2 * max_vertical + 2 * max_horizontal

    if num_slides > max_positions:
        raise ValueError("Too many slides to fit progress indicators on slide edges.")

    # Generate edge positions: right → bottom → left → top
    positions = []

    for i in range(max_vertical):
        x = width - margin - diameter
        y = margin + i * step
        positions.append((x, y))

    for i in range(max_horizontal):
        x = width - margin - diameter - i * step
        y = height - margin - diameter
        positions.append((x, y))

    for i in range(max_vertical):
        x = margin
        y = height - margin - diameter - i * step
        positions.append((x, y))

    for i in range(max_horizontal):
        x = margin + i * step
        y = margin
        positions.append((x, y))

    # Add indicators to each slide
    for slide_index, slide in enumerate(prs.slides):
        for i in range(num_slides):
            from pptx.enum.shapes import MSO_SHAPE

            # Position & size
            margin_left = Pt(20)
            margin_bottom = Pt(20)
            hourglass_height = Pt(100)
            hourglass_width = Pt(40)

            top = prs.slide_height - margin_bottom - hourglass_height
            center_y = top + hourglass_height // 2
            center_x = margin_left + hourglass_width // 2

            # Sand progress ratio
            progress = slide_index / (num_slides - 1) if num_slides > 1 else 1

            # Colors
            sand_color = RGBColor(255, 204, 0)
            glass_color = RGBColor(180, 180, 180)

            # -- Draw top sand (shrinking ellipse) --
            top_sand_height = int((1 - progress) * (hourglass_height // 2))
            top_sand = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                int(margin_left),
                int(center_y - top_sand_height),
                int(hourglass_width),
                int(top_sand_height)
            )
            top_sand.fill.solid()
            top_sand.fill.fore_color.rgb = sand_color
            top_sand.fill.fore_color.transparency = 0.1
            top_sand.line.fill.background()

            # --- Bottom bowl (glass outline) ---
            bottom_bowl_top = int(center_y)
            bottom_bowl_height = int(hourglass_height // 2)

            bottom_ellipse = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                int(margin_left),
                bottom_bowl_top,
                int(hourglass_width),
                bottom_bowl_height
            )
            bottom_ellipse.fill.background()
            bottom_ellipse.line.color.rgb = glass_color
            bottom_ellipse.line.width = Pt(1.25)

            # --- Bottom fill (fake sand, ellipse growing from bottom) ---
            # Simulate with an oval whose height increases and is vertically offset
            max_fill_height = bottom_bowl_height
            current_fill_height = int(progress * max_fill_height)

            # To center horizontally and place vertically aligned with bowl bottom
            fill_top = bottom_bowl_top + max_fill_height - current_fill_height

            sand_fill = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                int(margin_left),
                fill_top,
                int(hourglass_width),
                current_fill_height
            )
            sand_fill.fill.solid()
            sand_fill.fill.fore_color.rgb = sand_color
            sand_fill.fill.fore_color.transparency = 0.1
            sand_fill.line.fill.background()



            # -- Draw neck as a narrow rounded rectangle --
            neck = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(margin_left + hourglass_width // 3),
                int(center_y - Pt(2)),
                int(hourglass_width // 3),
                int(Pt(4))
            )
            neck.fill.solid()
            neck.fill.fore_color.rgb = sand_color
            neck.fill.fore_color.transparency = 0.3
            neck.line.fill.background()

            # -- Draw glass outline (optional) --
            glass = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                int(margin_left),
                int(top),
                int(hourglass_width),
                int(hourglass_height)
            )
            glass.fill.background()
            glass.line.color.rgb = glass_color
            glass.line.width = Pt(1.5)




    last_slide = prs.slides[-1]
    textbox = last_slide.shapes.add_textbox(
        left=Pt(30), top=height - Pt(80), width=width - Pt(60), height=Pt(40)
    )
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Total Slides: {num_slides}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    prs.save(output_path)
    print(f"Progress indicators and total count added to: {output_path}")

name = 'Title Lorem Ipsum'
add_wrapping_progress_circles(f'{name}.pptx', f'{name}_hourglass3.pptx')